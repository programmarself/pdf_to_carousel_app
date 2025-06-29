
import zipfile
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from docx import Document

def create_zip(images):
    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
        for idx, image in enumerate(images):
            img_byte_arr = BytesIO()
            image.save(img_byte_arr, format='PNG')
            zip_file.writestr(f'slide_{idx + 1}.png', img_byte_arr.getvalue())
    zip_buffer.seek(0)
    return zip_buffer

def create_ppt(images):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank_slide_layout)
        img_byte_arr = BytesIO()
        image.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)
        slide.shapes.add_picture(img_byte_arr, Inches(0), Inches(0), width=prs.slide_width)

    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

def create_word(images):
    doc = Document()
    for image in images:
        img_byte_arr = BytesIO()
        image.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)
        doc.add_picture(img_byte_arr, width=doc.sections[0].page_width - Inches(1))
        doc.add_page_break()

    word_buffer = BytesIO()
    doc.save(word_buffer)
    word_buffer.seek(0)
    return word_buffer
